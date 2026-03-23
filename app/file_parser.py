"""Extract text from PDF, PPTX, and image files for summary context."""

import os
from pathlib import Path


ALLOWED_EXTENSIONS = {".pdf", ".pptx", ".ppt", ".docx", ".doc", ".png", ".jpg", ".jpeg"}


def extract_text(filepath: str) -> str:
    """Extract text from a file based on its extension."""
    ext = Path(filepath).suffix.lower()
    if ext == ".pdf":
        return _extract_pdf(filepath)
    elif ext in (".pptx", ".ppt"):
        return _extract_pptx(filepath)
    elif ext in (".docx", ".doc"):
        return _extract_docx(filepath)
    elif ext in (".png", ".jpg", ".jpeg"):
        return _extract_image(filepath)
    return ""


def _extract_pdf(filepath: str) -> str:
    import fitz  # PyMuPDF
    text_parts = []
    with fitz.open(filepath) as doc:
        for page in doc:
            text_parts.append(page.get_text())
    return "\n".join(text_parts).strip()


def _extract_pptx(filepath: str) -> str:
    from pptx import Presentation
    text_parts = []
    prs = Presentation(filepath)
    for i, slide in enumerate(prs.slides, 1):
        slide_texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        slide_texts.append(t)
        if slide_texts:
            text_parts.append(f"[Diapositiva {i}]\n" + "\n".join(slide_texts))
    return "\n\n".join(text_parts).strip()


def _extract_docx(filepath: str) -> str:
    from docx import Document
    doc = Document(filepath)
    text_parts = []
    for para in doc.paragraphs:
        t = para.text.strip()
        if t:
            text_parts.append(t)
    return "\n".join(text_parts).strip()


def _extract_image(filepath: str) -> str:
    """Try OCR on image. Returns empty string if not available."""
    try:
        import fitz  # PyMuPDF can extract text from images via OCR
        doc = fitz.open(filepath)
        page = doc[0]
        text = page.get_text().strip()
        doc.close()
        return text
    except Exception:
        return ""


def extract_from_files(filepaths: list[str]) -> str:
    """Extract and combine text from multiple files."""
    parts = []
    for fp in filepaths:
        text = extract_text(fp)
        if text:
            name = Path(fp).name
            parts.append(f"[{name}]\n{text}")
    return "\n\n".join(parts)
